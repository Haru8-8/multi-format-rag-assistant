import os
import time
import numpy as np
import streamlit as st
from dotenv import load_dotenv
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from google import genai

# --- 1. アプリの初期設定 ---
load_dotenv()
st.set_page_config(page_title="ドキュメントAIアシスタント", layout="wide")
st.title("📄 資料読み込み型 AIチャット")

# APIクライアントの初期化
if "GOOGLE_API_KEY" in st.secrets:
    API_KEY = st.secrets["GOOGLE_API_KEY"]  # クラウド環境用
else:
    API_KEY = os.getenv("GOOGLE_API_KEY")  # ローカル開発（.env）用
client = genai.Client(api_key=API_KEY)

# --- 2. セッション状態（メモリ）の初期化 ---
# アプリ再読み込み時にも保持したい変数をここで定義します
if "chunks" not in st.session_state:
    st.session_state.chunks = []           # テキストチャンクとメタデータのリスト
if "embs" not in st.session_state:
    st.session_state.embs = []             # チャンクに対応するベクトルのリスト
if "processed_files" not in st.session_state:
    st.session_state.processed_files = []  # 処理済みファイル名のリスト（重複処理防止）
if "messages" not in st.session_state:
    st.session_state.messages = []         # チャットの会話履歴


# --- 3. コア機能（抽出・分割・ベクトル化） ---

def get_emb(text):
    """テキストをベクトル（数値の配列）に変換する"""
    res = client.models.embed_content(
        model="gemini-embedding-001",
        contents=text,
        config={'task_type': 'RETRIEVAL_DOCUMENT'}
    )
    return res.embeddings[0].values

def extract_text(uploaded_file):
    """各フォーマットに応じたテキスト抽出を行う"""
    file_ext = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""

    if file_ext == ".pdf":
        # PyMuPDFはファイルストリームから直接読み込める
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        for page in doc:
            text += page.get_text()
    elif file_ext == ".pptx":
        prs = Presentation(uploaded_file)
        for i, slide in enumerate(prs.slides):
            text += f"【スライド {i+1}】\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_ext == ".docx":
        doc = Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
        
    return text

def create_chunks(text, file_name, chunk_size=600, overlap=100):
    """抽出したテキストを意味の塊（チャンク）に分割し、メタデータを付与する"""
    chunks = []
    for i in range(0, len(text), chunk_size - overlap):
        chunk_text = text[i : i + chunk_size]
        if len(chunk_text.strip()) > 20: # 短すぎるゴミデータを除外
            chunks.append({
                "text": chunk_text,
                "source": file_name,
                "type": "normal"
            })
    return chunks

def generate_summary(text, file_name):
    """資料の全体像を把握するための「サマリーチャンク」を作成する"""
    # 巨大すぎるテキストはAPI制限に引っかかるため、先頭の一部を利用
    target_text = text[:5000] 
    
    prompt = f"以下の資料の内容を、主要なトピック、結論、特徴に焦点を当てて300文字程度で要約してください。:\n\n{target_text}"
    res = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    
    return {
        "text": f"【資料全体概要: {file_name}】\n{res.text}",
        "source": file_name,
        "type": "summary"
    }


# --- 4. RAG機能（検索・回答生成） ---

def rephrase_query(user_query, chat_history):
    """会話履歴とファイル一覧を元に、AIが検索しやすい具体的なクエリに書き換える"""
    if not chat_history:
        return user_query
    
    history_text = "\n".join([f"{m['role']}: {m['content']}" for m in chat_history[-3:]])
    file_list = ", ".join(st.session_state.processed_files)
    
    prompt = f"""
    これまでの会話の流れを読み、ユーザーの最新の質問を「検索しやすい具体的な文章」に書き換えてください。
    以下のファイルリストも参考にしてください。
    代名詞（それ、資料2など）があれば、具体的な名前や内容に置き換えてください。
    
    【ファイルリスト】: {file_list}
    【会話履歴】\n{history_text}
    【最新の質問】\n{user_query}
    
    回答は書き換えた文章のみを出力してください。
    """
    res = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    return res.text.strip()


# --- 5. UI構築（サイドバー：ナレッジベース構築） ---

with st.sidebar:
    st.header("ナレッジベース構築")
    uploaded_files = st.file_uploader("資料をアップロード (複数可)", type=["pdf", "pptx", "docx"], accept_multiple_files=True)
    
    for uploaded_file in uploaded_files:
        if uploaded_file.name not in st.session_state.processed_files:
            with st.spinner(f"{uploaded_file.name} を解析中..."):
                
                # 1. テキスト抽出
                raw_text = extract_text(uploaded_file)
                
                # 2. サマリーチャンクの作成（全体像の把握用）
                summary_chunk = generate_summary(raw_text, uploaded_file.name)
                
                # 3. 通常チャンクの作成と結合
                normal_chunks = create_chunks(raw_text, uploaded_file.name)
                all_chunks = [summary_chunk] + normal_chunks
                
                # 4. ベクトル化（RPM制限対策でsleepを入れる）
                new_embeddings = []
                progress_bar = st.progress(0)
                for i, item in enumerate(all_chunks):
                    new_embeddings.append(get_emb(item["text"]))
                    progress_bar.progress((i + 1) / len(all_chunks))
                    time.sleep(1.0) 

                # 5. セッション状態に保存
                st.session_state.chunks.extend(all_chunks)
                st.session_state.embs.extend(new_embeddings)
                st.session_state.processed_files.append(uploaded_file.name)
                
                st.success(f"{uploaded_file.name} の学習が完了しました！")


# --- 6. UI構築（メイン画面：チャット） ---

# 過去のメッセージを描画
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# ユーザーからの入力受け付け
if prompt := st.chat_input("資料の内容について質問してください"):
    if not st.session_state.chunks:
        st.warning("先にサイドバーから資料をアップロードしてください。")
    else:
        # ユーザーの質問を表示＆履歴に追加
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.spinner("AIが考え中..."):
            # 1. 質問の具体化
            refined_query = rephrase_query(prompt, st.session_state.messages)
            st.caption(f"🔍 内部検索クエリ: {refined_query}") # st.writeからcaptionに変更して控えめに表示

            # 2. ベクトル検索 (Top-8を取得)
            q_emb = get_emb(refined_query)
            sims = [np.dot(q_emb, d_emb) for d_emb in st.session_state.embs]
            top_idx = np.argsort(sims)[-8:][::-1]
            
            # 3. AIに渡すコンテキスト（参照テキスト）の作成
            context_parts = []
            for i in top_idx:
                item = st.session_state.chunks[i]
                context_parts.append(f"【資料名: {item['source']}】\n{item['text']}")
            context_text = "\n\n".join(context_parts)

            # 4. 回答の生成
            chat = client.chats.create(model="gemini-2.5-flash")
            final_prompt = f"""以下の「複数の資料」を参考にして答えてください。
            ユーザーが「資料1」や「パワポ」などと呼んでいる場合は、資料名から推測して適切に回答してください。
            回答の際、可能であれば「〇〇（資料名）によると〜」と出典を明記してください。
            資料に答えがない場合は、無理に答えず「資料には記載がありません」と伝えてください。

            【資料の抜粋】
            {context_text}

            【質問】
            {prompt}
            """
            
            response = chat.send_message(final_prompt)
            
            # 5. 回答と参照元の表示
            with st.chat_message("assistant"):
                st.markdown(response.text)
                with st.expander("📚 参照した箇所を確認"):
                    for c in context_parts:
                        st.info(c)
            
            # 履歴に追加
            st.session_state.messages.append({"role": "assistant", "content": response.text})