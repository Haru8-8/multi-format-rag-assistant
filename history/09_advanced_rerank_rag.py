import streamlit as st
from google import genai
import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import os
import time
import re
from dotenv import load_dotenv

load_dotenv()

# --- 設定 ---
st.set_page_config(page_title="最強RAGアシスタント", layout="wide")
st.title("📄 高精度マルチフォーマット RAG チャット")
st.caption("Multi-Query + HyDE + Rerank + Summary + Context Awareness")

API_KEY = st.secrets.get("GOOGLE_API_KEY") or os.getenv("GOOGLE_API_KEY")
client = genai.Client(api_key=API_KEY)

# --- 共通ユーティリティ ---

def get_emb(text):
    res = client.models.embed_content(
        model="gemini-embedding-001",
        contents=text,
        config={'task_type': 'RETRIEVAL_DOCUMENT'}
    )
    return res.embeddings[0].values

def extract_text(uploaded_file):
    file_ext = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""
    if file_ext == ".pdf":
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        text = "".join([page.get_text() for page in doc])
    elif file_ext == ".pptx":
        prs = Presentation(uploaded_file)
        for i, slide in enumerate(prs.slides):
            text += f"【スライド {i+1}】\n"
            text += "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]) + "\n"
    elif file_ext == ".docx":
        doc = Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
    return text

# --- 高度なロジック層 ---

def expand_query_with_hyde(user_query, chat_history):
    """ファイル名を考慮したMulti-QueryとHyDEの生成"""
    history_text = "\n".join([f"{m['role']}: {m['content']}" for m in chat_history[-3:]])
    file_list = ", ".join(st.session_state.processed_files)
    
    prompt = f"""あなたは高度な検索アシスタントです。
    現在アップロードされているファイル: {file_list}
    
    これらを踏まえ、ユーザーの質問に対し検索精度を最大化するための4点を出力してください。
    1-3. 異なる視点の検索クエリ（3つ）
    4. 質問に対する「仮説回答(HyDE)」
    
    質問: {user_query}
    履歴: {history_text}
    
    出力形式:
    Q1: [内容]
    Q2: [内容]
    Q3: [内容]
    HYDE: [内容]"""
    
    res = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    lines = res.text.strip().split('\n')
    queries = [user_query]
    hyde = ""
    for line in lines:
        if ":" in line:
            prefix, content = line.split(":", 1)
            content = content.strip()
            if not content: continue # 内容が空ならスキップ
            
            if prefix.strip() in ["Q1", "Q2", "Q3"]:
                queries.append(content)
            elif prefix.strip() == "HYDE":
                hyde = content
                
    # 空の文字列を除去してリストを作成
    result = [q for q in queries if q] 
    if hyde:
        result.append(hyde)
        
    return result

def rerank_chunks(query, chunks, chat_history, processed_files, top_n=5):
    """会話履歴とファイル名を考慮してチャンクを再ランキング"""
    if not chunks: return []
    
    # 直近の履歴とファイルリストをテキスト化
    history_text = "\n".join([f"{m['role']}: {m['content']}" for m in chat_history[-3:]])
    file_list = ", ".join(processed_files)
    
    chunks_text = ""
    for i, c in enumerate(chunks):
        # 出典（ファイル名）を明記することで、AIが識別しやすくする
        chunks_text += f"---\nID: {i}\n出典: {c['source']}\n内容: {c['text'][:400]}\n"
    
    prompt = f"""あなたは資料選別の専門家です。
    【現在読み込んでいるファイル】: {file_list}
    
    【これまでの会話】
    {history_text}
    
    【最新の質問】
    {query}
    
    上記の文脈とファイル名を最大限に考慮し、質問への回答に最も直結する資料を最大{top_n}個選び、関連性の高い順にIDを並べてください。
    特に「それら」や「共通点」といった言葉が指す対象を、会話履歴から正しく推測してください。
    
    【資料リスト】:
    {chunks_text}
    
    出力形式: IDのみをカンマ区切りで。例: 5, 2, 8"""
    
    try:
        res = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        ids = [int(x.strip()) for x in re.findall(r'\d+', res.text)]
        return [chunks[i] for i in ids if i < len(chunks)][:top_n]
    except:
        return chunks[:top_n]

# --- メインロジック ---

if "chunks" not in st.session_state:
    st.session_state.update({"chunks": [], "embs": [], "processed_files": [], "messages": []})

with st.sidebar:
    st.header("ナレッジベース構築")
    uploaded_files = st.file_uploader("資料アップロード", type=["pdf", "pptx", "docx"], accept_multiple_files=True)
    
    for f in uploaded_files:
        if f.name not in st.session_state.processed_files:
            # 1. テキスト抽出
            with st.spinner(f"{f.name} を読み込み中..."):
                raw_text = extract_text(f)
            
            # 2. 要約の生成
            with st.spinner(f"{f.name} の要約を作成中..."):
                summary_prompt = f"以下の資料の内容を300文字程度で簡潔に要約してください。:\n\n{raw_text[:5000]}"
                s_res = client.models.generate_content(model="gemini-2.5-flash-lite", contents=summary_prompt)
                
                # 要約を最初のチャンクとして登録
                doc_chunks = [{"text": f"【全体要約】\n{s_res.text}", "source": f.name, "type": "summary"}]
                
                # 通常チャンクの分割
                for i in range(0, len(raw_text), 500):
                    chunk_text = raw_text[i:i+600]
                    if len(chunk_text.strip()) > 20:
                        doc_chunks.append({"text": chunk_text, "source": f.name, "type": "normal"})
            
            # 3. ベクトル化（ここでプログレスバーを表示）
            st.write(f"📊 {f.name} のベクトル化を実行中...")
            p_bar = st.progress(0)
            
            for i, c in enumerate(doc_chunks):
                # ベクトル取得
                st.session_state.embs.append(get_emb(c["text"]))
                st.session_state.chunks.append(c)
                
                # プログレスバーの更新
                progress = (i + 1) / len(doc_chunks)
                p_bar.progress(progress)
                
                # API制限(Rate Limit)対策の待機
                time.sleep(0.2) 
            
            p_bar.empty() # 完了したらバーを消す
            st.session_state.processed_files.append(f.name)
            st.success(f"✅ {f.name} の学習が完了しました")

# チャットUI
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]): st.markdown(msg["content"])

if prompt := st.chat_input("質問をどうぞ"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"): st.markdown(prompt)

    if not st.session_state.chunks:
        st.warning("資料をアップロードしてください")
    else:
        with st.spinner("思考・検索・選別中..."):
            # Step 1: Query Expansion (Multi-Query + HyDE)
            targets = expand_query_with_hyde(prompt, st.session_state.messages)
            
            # Step 2: Retrieval (多めに回収)
            all_indices = set()
            for t in targets:
                if not t or not t.strip(): # 空の文字列を完全に除外
                    continue
                q_emb = get_emb(t)
                sims = [np.dot(q_emb, d_emb) for d_emb in st.session_state.embs]
                all_indices.update(np.argsort(sims)[-15:])
            
            candidates = [st.session_state.chunks[i] for i in all_indices]
            
            # Step 3: Rerank (Geminiによる精査、履歴とファイル名を渡す)
            refined_chunks = rerank_chunks(
                prompt, 
                candidates, 
                st.session_state.messages, 
                st.session_state.processed_files, 
                top_n=5
            )
            
            # Step 4: Final Generation
            context = "\n\n".join([f"【出典: {c['source']} ({c['type']})】\n{c['text']}" for c in refined_chunks])
            final_prompt = f"以下の資料を参考に質問に答えて。資料にないことは「不明」としてください。\n\n{context}\n\n質問: {prompt}"
            res = client.models.generate_content(model="gemini-2.5-flash", contents=final_prompt)
            
            with st.chat_message("assistant"):
                st.markdown(res.text)
                with st.expander("🔍 内部処理の詳細 (Rerank結果)"):
                    st.write("**生成クエリとHyDE:**", targets)
                    for i, c in enumerate(refined_chunks):
                        st.info(f"Rank {i+1}: {c['source']} ({c['type']})\n{c['text'][:200]}...")
            
            st.session_state.messages.append({"role": "assistant", "content": res.text})