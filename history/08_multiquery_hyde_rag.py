import streamlit as st
from google import genai
import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import os
import time
import re
from dotenv import load_dotenv

load_dotenv()

# --- 設定 ---
st.set_page_config(page_title="高度なドキュメントAIアシスタント", layout="wide")
st.title("📄 高精度マルチフォーマット RAG チャット")
st.caption("Multi-Query + HyDE 実装済みモデル")

# APIキーの設定（Secrets/env両対応）
if "GOOGLE_API_KEY" in st.secrets:
    API_KEY = st.secrets["GOOGLE_API_KEY"]  # クラウド環境用
else:
    API_KEY = os.getenv("GOOGLE_API_KEY")  # ローカル開発（.env）用
client = genai.Client(api_key=API_KEY)

# --- 共通ユーティリティ ---

def get_emb(text):
    """テキストをベクトル化する"""
    res = client.models.embed_content(
        model="gemini-embedding-001",
        contents=text,
        config={'task_type': 'RETRIEVAL_DOCUMENT'}
    )
    return res.embeddings[0].values

def extract_text(uploaded_file):
    """各ファイル形式からテキストを抽出"""
    file_ext = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""
    if file_ext == ".pdf":
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        text = "".join([page.get_text() for page in doc])
    elif file_ext == ".pptx":
        prs = Presentation(uploaded_file)
        for i, slide in enumerate(prs.slides):
            text += f"【スライド {i+1}】\n"
            text += "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]) + "\n"
    elif file_ext == ".docx":
        doc = Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
    return text

# --- クエリ拡張ロジック (Multi-Query + HyDE) ---

def expand_query_with_hyde(user_query, chat_history):
    """質問を多角化し、仮説回答(HyDE)を生成する"""
    history_text = "\n".join([f"{m['role']}: {m['content']}" for m in chat_history[-3:]])
    file_list = ", ".join(st.session_state.processed_files)
    
    prompt = f"""
    あなたは高度な検索アシスタントです。ユーザーの質問に対し、検索精度を最大化するために以下の4点を出力してください。
    
    1-3. 異なる視点からの検索クエリ（3つ）: 言葉の揺れや、関連するキーワードを含めてください。
    4. 質問に対する「仮説回答(HyDE)」: 資料内に存在しそうな具体的かつ技術的な回答文を1つ作成してください。

    【対象ファイル】: {file_list}
    【会話履歴】:
    {history_text}
    【ユーザーの質問】: {user_query}
    
    出力形式（厳守）:
    Q1: [クエリ1]
    Q2: [クエリ2]
    Q3: [クエリ3]
    HYDE: [仮説回答]
    """
    
    res = client.models.generate_content(model="gemini-2.5-flash-lite", contents=prompt)
    lines = res.text.strip().split('\n')
    
    expanded_queries = []
    hyde_text = ""
    for line in lines:
        if line.startswith(("Q1:", "Q2:", "Q3:")):
            expanded_queries.append(line.split(":", 1)[1].strip())
        elif line.startswith("HYDE:"):
            hyde_text = line.split(":", 1)[1].strip()
            
    # 元の質問、拡張クエリ、HyDEを混ぜたリストを返す
    return [user_query] + expanded_queries + [hyde_text]

# --- セッション初期化 ---
for key in ["chunks", "embs", "processed_files", "messages"]:
    if key not in st.session_state:
        st.session_state[key] = []

# --- サイドバー：ナレッジベース構築 ---
with st.sidebar:
    st.header("ナレッジベース構築")
    uploaded_files = st.file_uploader("資料アップロード", type=["pdf", "pptx", "docx"], accept_multiple_files=True)
    
    for uploaded_file in uploaded_files:
        if uploaded_file.name not in st.session_state.processed_files:
            with st.spinner(f"{uploaded_file.name} を解析中..."):
                raw_text = extract_text(uploaded_file)
                
                # 要約チャンクの作成
                summary_prompt = f"以下の資料を300文字程度で要約してください:\n\n{raw_text[:5000]}"
                s_res = client.models.generate_content(model="gemini-2.5-flash-lite", contents=summary_prompt)
                
                # チャンク分割
                new_chunks = [{"text": f"【概要: {uploaded_file.name}】\n{s_res.text}", "source": uploaded_file.name, "type": "summary"}]
                for i in range(0, len(raw_text), 500):
                    chunk = raw_text[i:i+600]
                    if len(chunk.strip()) > 20:
                        new_chunks.append({"text": chunk, "source": uploaded_file.name, "type": "normal"})
                
                # ベクトル化（429回避のためsleep）
                new_embs = []
                p_bar = st.progress(0)
                for i, c in enumerate(new_chunks):
                    new_embs.append(get_emb(c["text"]))
                    p_bar.progress((i + 1) / len(new_chunks))
                    time.sleep(0.5)
                
                st.session_state.chunks.extend(new_chunks)
                st.session_state.embs.extend(new_embs)
                st.session_state.processed_files.append(uploaded_file.name)
                st.success(f"{uploaded_file.name} 完了")

# --- メインチャット ---
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

if prompt := st.chat_input("質問をどうぞ"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    if not st.session_state.chunks:
        st.warning("資料をアップロードしてください")
    else:
        with st.spinner("思考中 (Multi-Query + HyDE)..."):
            # 1. クエリ拡張
            search_targets = expand_query_with_hyde(prompt, st.session_state.messages)
            
            with st.expander("🔍 生成された検索クエリと仮説回答"):
                for t in search_targets:
                    st.caption(t)

            # 2. 統合検索（複数クエリでそれぞれ検索）
            all_top_indices = set()
            for target_text in search_targets:
                q_emb = get_emb(target_text)
                sims = [np.dot(q_emb, d_emb) for d_emb in st.session_state.embs]
                # 各クエリごとに上位5つをピックアップ
                top_5 = np.argsort(sims)[-5:]
                all_top_indices.update(top_5)
            
            # 3. コンテキスト構築
            context_parts = []
            for i in all_top_indices:
                item = st.session_state.chunks[i]
                context_parts.append(f"【出典: {item['source']}】\n{item['text']}")
            context_text = "\n\n".join(context_parts)

            # 4. 回答生成
            final_prompt = f"""
            以下の複数の資料抜粋を参考に、ユーザーの質問に正確に答えてください。
            
            【資料抜粋】
            {context_text}
            
            【質問】
            {prompt}
            """
            res = client.models.generate_content(model="gemini-2.5-flash-lite", contents=final_prompt)
            
            with st.chat_message("assistant"):
                st.markdown(res.text)
                with st.expander("📚 参照した全チャンク"):
                    for cp in context_parts:
                        st.info(cp)
            
            st.session_state.messages.append({"role": "assistant", "content": res.text})