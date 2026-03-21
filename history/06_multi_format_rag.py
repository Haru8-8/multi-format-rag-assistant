import streamlit as st
from google import genai
import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation  # 追加
import os
import time
from dotenv import load_dotenv

load_dotenv()

# --- 設定 ---
st.set_page_config(page_title="ドキュメントAIアシスタント", layout="wide")
st.title("📄 資料読み込み型 AIチャット")

API_KEY = os.getenv("GOOGLE_API_KEY")
client = genai.Client(api_key=API_KEY)

def get_emb(text):
    res = client.models.embed_content(
        model="gemini-embedding-001",
        contents=text,
        config={'task_type': 'RETRIEVAL_DOCUMENT'}
    )
    return res.embeddings[0].values

# --- 抽出関数の拡張 ---

def extract_from_pptx(pptx_file):
    """PowerPointからテキストを抽出する"""
    prs = Presentation(pptx_file)
    text_list = []
    for i, slide in enumerate(prs.slides):
        slide_text = f"【スライド {i+1}】\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text_list.append(slide_text)
    return "\n".join(text_list)

def extract_and_chunk(uploaded_file, chunk_size=600, overlap=100):
    file_ext = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""

    if file_ext == ".pdf":
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        for page in doc:
            text += page.get_text()
    elif file_ext == ".pptx":
        text = extract_from_pptx(uploaded_file)
    
    # 共通のチャンク分割ロジック
    chunks = []
    for i in range(0, len(text), chunk_size - overlap):
        chunk = text[i : i + chunk_size]
        if len(chunk.strip()) > 20:
            chunks.append(chunk)
    return chunks

# --- セッション状態の初期化（アプリの冒頭で1回だけ実行） ---
if "chunks" not in st.session_state:
    st.session_state.chunks = []
if "embs" not in st.session_state:
    st.session_state.embs = []
if "processed_files" not in st.session_state:
    st.session_state.processed_files = []

# --- サイドバー：ファイルアップロード ---
with st.sidebar:
    st.header("ナレッジベース構築")
    uploaded_files = st.file_uploader("資料をアップロード (複数可)", type=["pdf", "pptx"], accept_multiple_files=True)
    
    for uploaded_file in uploaded_files:
        # まだ処理していないファイルだけを処理する
        if uploaded_file.name not in st.session_state.processed_files:
            with st.spinner(f"新着ファイル {uploaded_file.name} を解析・追加中..."):
                
                # 1. テキスト抽出 & チャンク作成
                new_chunks = extract_and_chunk(uploaded_file)
                
                # 2. ベクトル化 (RPM制限を考慮)
                new_embeddings = []
                progress_bar = st.progress(0)
                for i, c in enumerate(new_chunks):
                    new_embeddings.append(get_emb(c))
                    progress_bar.progress((i + 1) / len(new_chunks))
                    time.sleep(1.0)
                
                # 3. 既存のリストに「追加（結合）」する
                st.session_state.chunks.extend(new_chunks)
                st.session_state.embs.extend(new_embeddings)
                
                # 4. 処理済みリストに記録
                st.session_state.processed_files.append(uploaded_file.name)
                st.success(f"{uploaded_file.name} をナレッジに追加しました！")

# --- メインチャット画面 ---
if "messages" not in st.session_state:
    st.session_state.messages = []

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

if prompt := st.chat_input("資料の内容について質問してください"):
    if "chunks" not in st.session_state:
        st.warning("先にPDFファイルをアップロードしてください。")
    else:
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # 1. 質問のベクトル化
        q_emb = get_emb(prompt)
        
        # 2. 類似度計算 (Top-3)
        sims = [np.dot(q_emb, d_emb) for d_emb in st.session_state.embs]
        top_idx = np.argsort(sims)[-8:][::-1]
        relevant_chunks = [st.session_state.chunks[i] for i in top_idx]
        
        # 3. 回答生成
        context = "\n---\n".join(relevant_chunks)
        final_prompt = f"""以下の資料の一部を参考にして、質問に答えてください。
        資料に答えがない場合は、無理に答えず「資料には記載がありません」と伝えてください。

        【資料の内容】
        {context}

        【質問】
        {prompt}
        """
        
        response = client.models.generate_content(model="gemini-2.5-flash-lite", contents=final_prompt)
        
        with st.chat_message("assistant"):
            st.markdown(response.text)
            with st.expander("参照した箇所を確認"):
                for c in relevant_chunks:
                    st.info(c)
        
        st.session_state.messages.append({"role": "assistant", "content": response.text})